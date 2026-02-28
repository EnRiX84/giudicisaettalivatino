"""
Script per generare la presentazione PowerPoint per il Team Digitale
IISS "Giudici Saetta e Livatino"
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# Colori tema
BLU = RGBColor(0x1A, 0x3A, 0x6B)
BLU_SCURO = RGBColor(0x0F, 0x28, 0x47)
ORO = RGBColor(0xC9, 0xA8, 0x4C)
BIANCO = RGBColor(0xFF, 0xFF, 0xFF)
GRIGIO = RGBColor(0xF5, 0xF6, 0xFA)
TESTO = RGBColor(0x2D, 0x2D, 0x2D)
VERDE = RGBColor(0x19, 0x87, 0x54)
ROSSO = RGBColor(0xDC, 0x35, 0x45)

SCREENSHOTS_DIR = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'screenshots')

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def add_screenshot(slide, path, left, top, width, height=None):
    """Add screenshot with optional border"""
    if os.path.exists(path):
        if height:
            pic = slide.shapes.add_picture(path, Inches(left), Inches(top), Inches(width), Inches(height))
        else:
            pic = slide.shapes.add_picture(path, Inches(left), Inches(top), Inches(width))
        pic.line.color.rgb = RGBColor(0xE0, 0xE3, 0xEA)
        pic.line.width = Pt(2)
        return pic
    return None

def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_textbox(slide, left, top, width, height, text, font_size=18, color=TESTO, bold=False, alignment=PP_ALIGN.LEFT, font_name='Calibri'):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_bullet_list(slide, left, top, width, height, items, font_size=16, color=TESTO):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = 'Calibri'
        p.space_after = Pt(6)
        p.level = 0
    return txBox

def add_shape_bg(slide, left, top, width, height, color, alpha=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    if alpha is not None:
        shape.fill.fore_color.brightness = 0
    return shape

def add_card(slide, left, top, width, height, title, body, title_color=BLU, bg_color=BIANCO):
    # Card background
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    card.fill.solid()
    card.fill.fore_color.rgb = bg_color
    card.line.color.rgb = RGBColor(0xE0, 0xE3, 0xEA)
    card.line.width = Pt(1)
    card.shadow.inherit = False
    # Title
    add_textbox(slide, left + 0.3, top + 0.2, width - 0.6, 0.5, title, 16, title_color, True)
    # Body
    add_textbox(slide, left + 0.3, top + 0.7, width - 0.6, height - 1, body, 13, TESTO)
    return card

# =============================================
# SLIDE 1 - TITOLO
# =============================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_shape_bg(slide, 0, 0, 13.333, 7.5, BLU_SCURO)
# Linea oro decorativa
add_shape_bg(slide, 0, 3.2, 13.333, 0.05, ORO)

add_textbox(slide, 1, 1.5, 11.333, 1.2,
    'IISS "Giudici Saetta e Livatino"', 44, BIANCO, True, PP_ALIGN.CENTER)
add_textbox(slide, 1, 2.5, 11.333, 0.6,
    'Ravanusa e Campobello di Licata', 22, ORO, False, PP_ALIGN.CENTER)

add_textbox(slide, 1, 3.8, 11.333, 1.2,
    'Proposta di adozione del nuovo sito web istituzionale', 32, BIANCO, True, PP_ALIGN.CENTER)
add_textbox(slide, 1, 5.0, 11.333, 0.6,
    'Presentazione per il Team Digitale - Febbraio 2026', 18, RGBColor(0xAA, 0xAA, 0xCC), False, PP_ALIGN.CENTER)

add_textbox(slide, 1, 5.8, 11.333, 0.5,
    'Conforme alle Linee Guida AgID e WCAG 2.1 AA', 16, ORO, True, PP_ALIGN.CENTER)

add_textbox(slide, 1, 6.5, 11.333, 0.5,
    'A cura del Prof. Enrico Maria Caruso - Docente di Informatica',
    15, RGBColor(0xBB, 0xBB, 0xDD), False, PP_ALIGN.CENTER)

# =============================================
# SLIDE 2 - SITUAZIONE ATTUALE
# =============================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BIANCO)
add_shape_bg(slide, 0, 0, 13.333, 1.2, BLU)
add_textbox(slide, 0.8, 0.25, 11.733, 0.7, 'Il sito attuale: le criticità', 32, BIANCO, True, PP_ALIGN.LEFT)

items = [
    'Template Joomla JSN Pixel 2 risalente ai primi anni 2010',
    'Non responsive: inutilizzabile su smartphone e tablet',
    'Caricamento lento (3-6 secondi) per plugin e query database',
    'Vulnerabilità di sicurezza: richiede aggiornamenti costanti di Joomla e PHP',
    'Menu confuso con voci ridondanti, dropdown non funzionanti su mobile',
    'Non conforme ai requisiti di accessibilità AgID (Legge 4/2004)',
    'Documenti difficili da reperire, struttura dei contenuti poco intuitiva',
    'Costi ricorrenti: hosting PHP/MySQL 100-300+ euro/anno + manutenzione',
]

for i, item in enumerate(items):
    y = 1.6 + i * 0.65
    # Icona X rossa
    add_textbox(slide, 0.8, y, 0.5, 0.5, 'X', 18, ROSSO, True)
    add_textbox(slide, 1.4, y, 11, 0.55, item, 17, TESTO)

# =============================================
# SLIDE 3 - IL NUOVO SITO
# =============================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BIANCO)
add_shape_bg(slide, 0, 0, 13.333, 1.2, BLU)
add_textbox(slide, 0.8, 0.25, 11.733, 0.7, 'Il nuovo sito: moderno, veloce, sicuro', 32, BIANCO, True)

items = [
    'Design moderno e accattivante con palette istituzionale blu/oro',
    'Mobile-first: navigazione perfetta su smartphone e tablet',
    'Caricamento istantaneo (<1 secondo) - sito statico senza database',
    'Sicurezza intrinseca: zero vulnerabilità server-side, nessun CMS da aggiornare',
    'Navigazione intuitiva con navbar sticky e menu ben organizzati',
    'Conforme WCAG 2.1 AA e Linee Guida AgID sull\'accessibilità',
    'Card grid per notizie, sezioni tematiche, documenti facilmente reperibili',
    'Hosting gratuito (GitHub Pages, Netlify) - zero costi ricorrenti',
]

for i, item in enumerate(items):
    y = 1.6 + i * 0.65
    add_textbox(slide, 0.8, y, 0.5, 0.5, '\u2713', 20, VERDE, True)
    add_textbox(slide, 1.4, y, 11, 0.55, item, 17, TESTO)

# =============================================
# SLIDE 4 - CONFRONTO VISIVO
# =============================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, GRIGIO)
add_shape_bg(slide, 0, 0, 13.333, 1.2, BLU)
add_textbox(slide, 0.8, 0.25, 11.733, 0.7, 'Confronto: prima e dopo', 32, BIANCO, True)

# Tabella confronto
headers = ['Aspetto', 'Sito Attuale (Joomla)', 'Nuovo Sito (Statico)']
rows = [
    ['Design', 'Datato anni 2010, piatto', 'Moderno, palette blu/oro, ombre morbide'],
    ['Mobile', 'Non responsive', 'Mobile-first, touch-friendly'],
    ['Performance', '3-6 secondi caricamento', '<1 secondo caricamento'],
    ['Sicurezza', 'Vulnerabilità frequenti', 'Zero vulnerabilità server-side'],
    ['Accessibilità', 'Non conforme AgID', 'WCAG 2.1 AA conforme'],
    ['Manutenzione', 'Aggiornamenti costanti', 'Zero manutenzione server'],
    ['Costo annuo', '100-300+ euro', '0 euro'],
]

col_widths = [2.5, 4.5, 4.5]
start_x = 0.9
start_y = 1.5
row_h = 0.6

# Headers
for j, h in enumerate(headers):
    x = start_x + sum(col_widths[:j])
    cell = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(start_y), Inches(col_widths[j]), Inches(row_h))
    cell.fill.solid()
    cell.fill.fore_color.rgb = BLU
    cell.line.color.rgb = BLU
    tf = cell.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].text = h
    tf.paragraphs[0].font.size = Pt(14)
    tf.paragraphs[0].font.color.rgb = BIANCO
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.name = 'Calibri'
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

# Data rows
for i, row in enumerate(rows):
    y = start_y + (i + 1) * row_h
    bg = BIANCO if i % 2 == 0 else GRIGIO
    for j, val in enumerate(row):
        x = start_x + sum(col_widths[:j])
        cell = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(col_widths[j]), Inches(row_h))
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg
        cell.line.color.rgb = RGBColor(0xE0, 0xE3, 0xEA)
        cell.line.width = Pt(0.5)
        tf = cell.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].text = val
        tf.paragraphs[0].font.size = Pt(13)
        tf.paragraphs[0].font.name = 'Calibri'
        if j == 0:
            tf.paragraphs[0].font.bold = True
            tf.paragraphs[0].font.color.rgb = BLU
        elif j == 1:
            tf.paragraphs[0].font.color.rgb = ROSSO
        else:
            tf.paragraphs[0].font.color.rgb = VERDE
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER

# =============================================
# SLIDE 5 - CONFRONTO SCREENSHOT DESKTOP
# =============================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, GRIGIO)
add_shape_bg(slide, 0, 0, 13.333, 1.2, BLU)
add_shape_bg(slide, 0, 1.15, 13.333, 0.05, ORO)
add_textbox(slide, 0.8, 0.25, 11.733, 0.7, 'Confronto visivo: Homepage Desktop', 32, BIANCO, True)

# Prima - label
add_shape_bg(slide, 0.6, 1.5, 5.9, 0.5, ROSSO)
add_textbox(slide, 0.6, 1.5, 5.9, 0.5, 'PRIMA', 18, BIANCO, True, PP_ALIGN.CENTER)

# Dopo - label
add_shape_bg(slide, 6.8, 1.5, 5.9, 0.5, VERDE)
add_textbox(slide, 6.8, 1.5, 5.9, 0.5, 'DOPO', 18, BIANCO, True, PP_ALIGN.CENTER)

# Screenshots
add_screenshot(slide, os.path.join(SCREENSHOTS_DIR, 'vecchio_homepage.png'), 0.6, 2.1, 5.9, 4.5)
add_screenshot(slide, os.path.join(SCREENSHOTS_DIR, 'nuovo_homepage.png'), 6.8, 2.1, 5.9, 4.5)

add_textbox(slide, 0.6, 6.7, 5.9, 0.4, 'Joomla JSN Pixel 2 - Layout datato', 12, ROSSO, False, PP_ALIGN.CENTER)
add_textbox(slide, 6.8, 6.7, 5.9, 0.4, 'Design moderno - Palette istituzionale blu/oro', 12, VERDE, False, PP_ALIGN.CENTER)

# =============================================
# SLIDE 6 - CONFRONTO SCREENSHOT MOBILE
# =============================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, GRIGIO)
add_shape_bg(slide, 0, 0, 13.333, 1.2, BLU)
add_shape_bg(slide, 0, 1.15, 13.333, 0.05, ORO)
add_textbox(slide, 0.8, 0.25, 11.733, 0.7, 'Confronto visivo: Versione Mobile', 32, BIANCO, True)

# Prima - label
add_shape_bg(slide, 1.5, 1.5, 4.5, 0.5, ROSSO)
add_textbox(slide, 1.5, 1.5, 4.5, 0.5, 'PRIMA (non responsive)', 16, BIANCO, True, PP_ALIGN.CENTER)

# Dopo - label
add_shape_bg(slide, 7.5, 1.5, 4.5, 0.5, VERDE)
add_textbox(slide, 7.5, 1.5, 4.5, 0.5, 'DOPO (mobile-first)', 16, BIANCO, True, PP_ALIGN.CENTER)

# Screenshots mobile (centrati, piu' stretti)
add_screenshot(slide, os.path.join(SCREENSHOTS_DIR, 'vecchio_mobile.png'), 2.0, 2.2, 3.5, 5.0)
add_screenshot(slide, os.path.join(SCREENSHOTS_DIR, 'nuovo_mobile.png'), 8.0, 2.2, 3.5, 5.0)

# =============================================
# SLIDE 7 - DETTAGLIO NUOVO SITO
# =============================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, GRIGIO)
add_shape_bg(slide, 0, 0, 13.333, 1.2, BLU)
add_shape_bg(slide, 0, 1.15, 13.333, 0.05, ORO)
add_textbox(slide, 0.8, 0.25, 11.733, 0.7, 'Il nuovo sito nel dettaglio', 32, BIANCO, True)

# Screenshot notizie
add_screenshot(slide, os.path.join(SCREENSHOTS_DIR, 'nuovo_notizie.png'), 0.6, 1.5, 6.0, 3.8)
add_textbox(slide, 0.6, 5.4, 6.0, 0.4, 'Sezione Notizie - Card con categorie e toggle', 13, BLU, True, PP_ALIGN.CENTER)

# Screenshot offerta
add_screenshot(slide, os.path.join(SCREENSHOTS_DIR, 'nuovo_offerta.png'), 6.8, 1.5, 6.0, 3.8)
add_textbox(slide, 6.8, 5.4, 6.0, 0.4, 'Offerta Formativa - Card con gradiente per indirizzo', 13, BLU, True, PP_ALIGN.CENTER)

# Punti chiave sotto
add_shape_bg(slide, 0.6, 5.9, 12.133, 1.2, BIANCO)
key_points = [
    '\u2022 23 notizie organizzate in card con 6 categorie colorate',
    '\u2022 4 indirizzi di studio con card gradient distinte per colore',
    '\u2022 Toggle "Mostra tutte" per le notizie nascoste (17 su 23)',
    '\u2022 Animazioni fade-in on scroll per un\'esperienza moderna',
]
for i, point in enumerate(key_points):
    add_textbox(slide, 0.8, 6.0 + i * 0.28, 11.733, 0.3, point, 12, TESTO)

# =============================================
# SLIDE 8 - FUNZIONALITA'
# =============================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BIANCO)
add_shape_bg(slide, 0, 0, 13.333, 1.2, BLU)
add_textbox(slide, 0.8, 0.25, 11.733, 0.7, 'Funzionalità del nuovo sito', 32, BIANCO, True)

cards_data = [
    ('25 pagine interne', 'Organigramma, regolamenti,\nmodulistica, contatti, sedi,\ncalendario, orario lezioni...'),
    ('23 notizie e avvisi', 'Card con layout a griglia,\ncategorie colorate, toggle\nmostra/nascondi integrato'),
    ('Offerta formativa', '4 indirizzi di studio con\nschede dedicate: Scientifico,\nScienze Umane, AFM, Biomedica'),
    ('7 sezioni progetti', 'Erasmus+, PNRR, PON FESR,\nPCTO, Cambridge, STEAM,\nWeCanJob'),
    ('15 banner istituzionali', 'Amm. Trasparente, Albo\nPretorio, Accesso Civico,\nPNRR, Erasmus+...'),
    ('Pannello admin', 'Gestione notizie via browser,\npubblicazione con un click,\nbackup automatico'),
]

for i, (title, body) in enumerate(cards_data):
    col = i % 3
    row = i // 3
    x = 0.6 + col * 4.2
    y = 1.6 + row * 2.8
    add_card(slide, x, y, 3.8, 2.4, title, body)

# =============================================
# SLIDE 6 - ACCESSIBILITA' WCAG
# =============================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BIANCO)
add_shape_bg(slide, 0, 0, 13.333, 1.2, BLU)
add_textbox(slide, 0.8, 0.25, 11.733, 0.7, 'Accessibilità: conformità WCAG 2.1 AA', 32, BIANCO, True)

add_textbox(slide, 0.8, 1.4, 12, 0.5,
    'Il sito è conforme alla Legge 4/2004 (Legge Stanca) e alle Linee Guida AgID', 18, TESTO, True)

left_items = [
    '\u2713  Skip link "Vai al contenuto principale"',
    '\u2713  Struttura semantica HTML5 (main, nav, footer)',
    '\u2713  Attributi ARIA su menu, dropdown, toggle',
    '\u2713  aria-expanded sui controlli espandibili',
    '\u2713  aria-hidden sulle icone decorative (SVG)',
    '\u2713  Navigazione completa da tastiera',
    '\u2713  Indicatori di focus visibili (outline oro 3px)',
]

right_items = [
    '\u2713  Contrasti >= 4.5:1 (testo) e >= 3:1 (UI)',
    '\u2713  Design responsive mobile-first',
    '\u2713  prefers-reduced-motion rispettato',
    '\u2713  Testi alternativi su immagini',
    '\u2713  Avviso "si apre in nuova finestra" per link esterni',
    '\u2713  Dichiarazione di accessibilità nel footer',
    '\u2713  Meccanismo di feedback per segnalazioni',
]

for i, item in enumerate(left_items):
    add_textbox(slide, 0.8, 2.1 + i * 0.6, 6, 0.5, item, 15, VERDE, False)

for i, item in enumerate(right_items):
    add_textbox(slide, 6.8, 2.1 + i * 0.6, 6, 0.5, item, 15, VERDE, False)

add_shape_bg(slide, 0.6, 6.3, 12.133, 0.7, RGBColor(0xE3, 0xF0, 0xFF))
add_textbox(slide, 0.8, 6.35, 11.733, 0.6,
    'Dichiarazione di Accessibilità compilabile su form.agid.gov.it entro il 23 settembre di ogni anno',
    14, BLU, True, PP_ALIGN.CENTER)

# =============================================
# SLIDE 7 - USABILITA' AgID
# =============================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BIANCO)
add_shape_bg(slide, 0, 0, 13.333, 1.2, BLU)
add_textbox(slide, 0.8, 0.25, 11.733, 0.7, 'Usabilità: Linee Guida AgID per i siti PA', 32, BIANCO, True)

usability_cards = [
    ('Mobile-First', 'Interfaccia ottimizzata per\nsmartphone e tablet.\nBreakpoint a 480px, 768px, 1024px.'),
    ('Navigazione intuitiva', 'Navbar sticky, hamburger menu,\ndropdown organizzati per\ncategoria. Max 2 click.'),
    ('Contenuti accessibili', 'Card grid per notizie, doc-list\nper documenti, tabelle\nper orari. Tutto a portata.'),
    ('Tipografia leggibile', 'Font Inter (corpo) e Playfair\nDisplay (titoli). Line-height\n1.6, contrasti elevati.'),
    ('Performance', 'Sito statico: caricamento\n<1 secondo. Nessun database,\nnessun framework pesante.'),
    ('Coerenza visiva', 'Palette istituzionale uniforme,\ncomponenti riutilizzabili,\ndesign system coerente.'),
]

for i, (title, body) in enumerate(usability_cards):
    col = i % 3
    row = i // 3
    x = 0.6 + col * 4.2
    y = 1.6 + row * 2.8
    add_card(slide, x, y, 3.8, 2.4, title, body)

# =============================================
# SLIDE 8 - PANNELLO ADMIN
# =============================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BIANCO)
add_shape_bg(slide, 0, 0, 13.333, 1.2, BLU)
add_textbox(slide, 0.8, 0.25, 11.733, 0.7, 'Gestione autonoma: il pannello admin', 32, BIANCO, True)

add_textbox(slide, 0.8, 1.5, 12, 0.8,
    'La segreteria gestisce le notizie in totale autonomia, senza competenze tecniche.', 20, TESTO, False)

steps = [
    ('1', 'Doppio click', 'La segretaria fa doppio click\nsu "Avvia Gestione Sito.bat"'),
    ('2', 'Browser si apre', 'Si apre automaticamente il\npannello admin nel browser'),
    ('3', 'Compila il form', 'Inserisce titolo, categoria,\ndata, descrizione, allegati'),
    ('4', 'Pubblica', 'Clicca "Pubblica sul sito"\ne le notizie sono online'),
]

for i, (num, title, desc) in enumerate(steps):
    x = 0.6 + i * 3.2
    # Cerchio numerato
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + 0.9), Inches(2.8), Inches(0.8), Inches(0.8))
    circle.fill.solid()
    circle.fill.fore_color.rgb = ORO
    circle.line.fill.background()
    tf = circle.text_frame
    tf.paragraphs[0].text = num
    tf.paragraphs[0].font.size = Pt(24)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = BLU_SCURO
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    add_textbox(slide, x, 3.8, 3, 0.5, title, 18, BLU, True, PP_ALIGN.CENTER)
    add_textbox(slide, x, 4.3, 3, 1.2, desc, 14, TESTO, False, PP_ALIGN.CENTER)

add_shape_bg(slide, 0.6, 5.8, 12.133, 1.2, GRIGIO)
add_textbox(slide, 0.8, 5.9, 11.733, 0.5,
    'Caratteristiche aggiuntive:', 16, BLU, True)
add_textbox(slide, 0.8, 6.3, 11.733, 0.6,
    '\u2022 Backup automatico ad ogni salvataggio   \u2022 Riordinamento notizie drag & drop   \u2022 Eliminazione con conferma   \u2022 Server locale leggero (Python)', 14, TESTO)

# =============================================
# SLIDE - AUTENTICAZIONE E GESTIONE UTENTI (NUOVO v3)
# =============================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BIANCO)
add_shape_bg(slide, 0, 0, 13.333, 1.2, BLU)
add_shape_bg(slide, 0, 1.15, 13.333, 0.05, ORO)
add_textbox(slide, 0.8, 0.25, 11.733, 0.7, 'Sicurezza: Autenticazione e Gestione Utenti', 32, BIANCO, True)

add_textbox(slide, 0.8, 1.5, 12, 0.5,
    'Il pannello admin è protetto da un sistema di login con ruoli e gestione utenti.', 18, TESTO, True)

# Colonna sinistra - Login e sicurezza
add_card(slide, 0.6, 2.2, 5.9, 4.5, 'Login e Sicurezza', '')
login_items = [
    '\u2713  Schermata di login con username e password',
    '\u2713  Password PBKDF2-SHA256, 100.000 iterazioni',
    '\u2713  Salt random per utente (constant-time compare)',
    '\u2713  Cookie HttpOnly + SameSite=Strict (no XSS/CSRF)',
    '\u2713  Sessioni con scadenza automatica (8 ore)',
    '\u2713  Cambio password con verifica della vecchia',
    '\u2713  Avviso password predefinita al primo avvio',
]
for i, item in enumerate(login_items):
    add_textbox(slide, 0.8, 2.9 + i * 0.48, 5.5, 0.45, item, 13, VERDE)

# Colonna destra - Ruoli e gestione
add_card(slide, 6.8, 2.2, 5.9, 4.5, 'Ruoli e Gestione Utenti', '')
role_items = [
    '\u2022  Due ruoli: Admin ed Editor',
    '\u2022  Admin: gestione completa + tab Utenti',
    '\u2022  Editor: gestione notizie e pagine',
    '\u2022  Creazione/eliminazione utenti (solo admin)',
    '\u2022  Cambio ruolo inline dalla tabella utenti',
    '\u2022  Protezione ultimo admin (non eliminabile)',
    '\u2022  Protezione auto-eliminazione',
]
for i, item in enumerate(role_items):
    add_textbox(slide, 7.0, 2.9 + i * 0.48, 5.5, 0.45, item, 13, TESTO)

add_shape_bg(slide, 0.6, 6.9, 12.133, 0.5, RGBColor(0xE3, 0xF0, 0xFF))
add_textbox(slide, 0.8, 6.93, 11.733, 0.4,
    'Nessuna dipendenza esterna: tutto integrato nel server Python locale. Utenti salvati in data/users.json.',
    13, BLU, True, PP_ALIGN.CENTER)

# =============================================
# SLIDE - COOKIE BANNER GDPR (NUOVO v3)
# =============================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BIANCO)
add_shape_bg(slide, 0, 0, 13.333, 1.2, BLU)
add_shape_bg(slide, 0, 1.15, 13.333, 0.05, ORO)
add_textbox(slide, 0.8, 0.25, 11.733, 0.7, 'Privacy: Cookie Banner conforme al GDPR', 32, BIANCO, True)

add_textbox(slide, 0.8, 1.5, 12, 0.5,
    'Banner informativo sui cookie conforme al Regolamento UE 2016/679 (GDPR) e alle linee guida del Garante Privacy.', 18, TESTO, True)

# Colonna sinistra - Caratteristiche banner
add_card(slide, 0.6, 2.3, 5.9, 2.5, 'Il Cookie Banner', '')
banner_items = [
    '\u2713  Appare alla prima visita su tutte le pagine',
    '\u2713  Tre opzioni: Accetta / Rifiuta / Personalizza',
    '\u2713  Pannello dettagli con toggle per categoria',
    '\u2713  Link "Preferenze cookie" nel footer del sito',
    '\u2713  Consenso salvato in localStorage (no cookie aggiuntivi)',
]
for i, item in enumerate(banner_items):
    add_textbox(slide, 0.8, 3.0 + i * 0.38, 5.5, 0.35, item, 12, VERDE)

# Colonna destra - Categorie cookie
add_card(slide, 6.8, 2.3, 5.9, 2.5, 'Categorie di Cookie', '')
cat_items = [
    '\u2022  Tecnici necessari: sessione admin (HttpOnly)',
    '    Sempre attivi, non disattivabili dall\'utente',
    '\u2022  Servizi terze parti: Google Fonts (tipografia)',
    '    Disattivabili dall\'utente tramite il pannello',
    '\u2022  Nessun cookie di profilazione o pubblicitario',
]
for i, item in enumerate(cat_items):
    add_textbox(slide, 7.0, 3.0 + i * 0.38, 5.5, 0.35, item, 12, TESTO)

# Conformita'
add_card(slide, 0.6, 5.0, 12.133, 2.0, 'Conformità normativa', '')
compliance_items = [
    '\u2713  GDPR (Reg. UE 2016/679) - Consenso informato e granulare per cookie non necessari',
    '\u2713  Codice Privacy italiano (D.lgs. 196/2003 aggiornato) - Informativa chiara e accessibile',
    '\u2713  Linee guida Garante Privacy 10 giugno 2021 - Banner con "Accetta", "Rifiuta" e "Personalizza"',
    '\u2713  Accessibilità WCAG 2.1 AA - Banner navigabile da tastiera, focus visibili, aria-label',
]
for i, item in enumerate(compliance_items):
    add_textbox(slide, 0.8, 5.7 + i * 0.35, 11.733, 0.33, item, 12, VERDE)

# =============================================
# SLIDE 9 - CARATTERISTICHE PECULIARI
# =============================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BIANCO)
add_shape_bg(slide, 0, 0, 13.333, 1.2, BLU)
add_shape_bg(slide, 0, 1.15, 13.333, 0.05, ORO)
add_textbox(slide, 0.8, 0.25, 11.733, 0.7, 'Caratteristiche peculiari del nuovo sito', 32, BIANCO, True)

peculiar = [
    ('\u2b50', 'Anno scolastico auto-aggiornante',
     'Si aggiorna automaticamente ogni 1 settembre\ntramite JavaScript, senza intervento manuale.'),
    ('\u26a1', 'Sito statico ultra-veloce',
     'Nessun database, nessun PHP, nessun server\ndi backend. Solo file HTML/CSS/JS puri.'),
    ('\u2699', 'Pannello admin integrato',
     'La segreteria pubblica notizie con un click.\nNessuna competenza tecnica richiesta.'),
    ('\u26c4', 'Backup automatico',
     'Ogni salvataggio crea una copia di sicurezza.\nImpossibile perdere dati accidentalmente.'),
    ('\u2714', '23 notizie + 25 pagine + 15 banner',
     'Tutti i contenuti del sito originale preservati\ne riorganizzati in modo moderno e intuitivo.'),
    ('\u2615', 'Server locale con un click',
     'Doppio click su un file .bat per avviare\nil server di gestione. Zero installazioni.'),
    ('\u267b', 'Link PDF tutti verificati',
     'Ogni link a documenti PDF è stato verificato\ne punta al dominio originale della scuola.'),
    ('\u2666', 'Copiabile su chiavetta USB',
     'L\'intero sito è una singola cartella.\nBackup = copia su chiavetta. Migrazione = copia.'),
]

for i, (icon, title, desc) in enumerate(peculiar):
    col = i % 4
    row = i // 4
    x = 0.4 + col * 3.25
    y = 1.5 + row * 2.8
    # Card background
    card_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(3.0), Inches(2.4))
    card_bg.fill.solid()
    card_bg.fill.fore_color.rgb = GRIGIO
    card_bg.line.color.rgb = RGBColor(0xE0, 0xE3, 0xEA)
    card_bg.line.width = Pt(1)
    # Barra oro in alto
    add_shape_bg(slide, x, y, 3.0, 0.06, ORO)
    add_textbox(slide, x + 0.2, y + 0.2, 0.5, 0.5, icon, 22, ORO, True)
    add_textbox(slide, x + 0.2, y + 0.7, 2.6, 0.5, title, 13, BLU, True)
    add_textbox(slide, x + 0.2, y + 1.2, 2.6, 1.1, desc, 11, TESTO)

# =============================================
# SLIDE 10 - DEPLOY DETTAGLIATO
# =============================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BIANCO)
add_shape_bg(slide, 0, 0, 13.333, 1.2, BLU)
add_shape_bg(slide, 0, 1.15, 13.333, 0.05, ORO)
add_textbox(slide, 0.8, 0.25, 11.733, 0.7, 'Come effettuare il deploy', 32, BIANCO, True)

# Opzione 1
add_shape_bg(slide, 0.6, 1.5, 5.9, 5.5, GRIGIO)
add_shape_bg(slide, 0.6, 1.5, 5.9, 0.06, ORO)
add_textbox(slide, 0.8, 1.7, 5.5, 0.5, 'OPZIONE A: Hosting esistente', 20, BLU, True)
add_textbox(slide, 0.8, 2.2, 5.5, 0.5, 'Pubblicare sullo STESSO server attuale', 14, ORO, True)

deploy_a_steps = [
    ('1', 'Accedere al pannello di controllo hosting (cPanel)'),
    ('2', 'Aprire il File Manager nella document root'),
    ('3', 'Spostare/rinominare la cartella Joomla attuale come backup'),
    ('4', 'Caricare TUTTI i file del nuovo sito (upload ZIP + estrai)'),
    ('5', 'Verificare che index.html sia nella root'),
    ('6', 'Il sito è online! Stessa URL, nessuna modifica DNS'),
]

for i, (num, step) in enumerate(deploy_a_steps):
    y = 2.9 + i * 0.6
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.9), Inches(y), Inches(0.4), Inches(0.4))
    circle.fill.solid()
    circle.fill.fore_color.rgb = BLU
    circle.line.fill.background()
    tf = circle.text_frame
    tf.paragraphs[0].text = num
    tf.paragraphs[0].font.size = Pt(12)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = BIANCO
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    add_textbox(slide, 1.5, y, 4.8, 0.4, step, 12, TESTO)

# Opzione 2
add_shape_bg(slide, 6.8, 1.5, 5.9, 5.5, GRIGIO)
add_shape_bg(slide, 6.8, 1.5, 5.9, 0.06, ORO)
add_textbox(slide, 7.0, 1.7, 5.5, 0.5, 'OPZIONE B: Hosting gratuito', 20, BLU, True)
add_textbox(slide, 7.0, 2.2, 5.5, 0.5, 'GitHub Pages / Netlify / Cloudflare Pages', 14, ORO, True)

deploy_b_steps = [
    ('1', 'Creare account gratuito su GitHub'),
    ('2', 'Creare un repository e caricare i file del sito'),
    ('3', 'Attivare GitHub Pages nelle impostazioni del repo'),
    ('4', 'Il sito è online su nomescuola.github.io'),
    ('5', 'Configurare dominio personalizzato (opzionale):\n     saettalivatinoravanusa.edu.it > GitHub Pages'),
    ('6', 'HTTPS automatico e gratuito incluso'),
]

for i, (num, step) in enumerate(deploy_b_steps):
    y = 2.9 + i * 0.6
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(7.1), Inches(y), Inches(0.4), Inches(0.4))
    circle.fill.solid()
    circle.fill.fore_color.rgb = VERDE
    circle.line.fill.background()
    tf = circle.text_frame
    tf.paragraphs[0].text = num
    tf.paragraphs[0].font.size = Pt(12)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = BIANCO
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    add_textbox(slide, 7.7, y, 4.8, 0.55, step, 12, TESTO)

# Nota bene
add_shape_bg(slide, 0.6, 7.0, 12.133, 0.35, RGBColor(0xE3, 0xF0, 0xFF))
add_textbox(slide, 0.8, 7.02, 11.733, 0.3,
    'Requisiti tecnici: NESSUNO. Il sito non richiede database, PHP, o configurazioni server. Solo file statici.',
    12, BLU, True, PP_ALIGN.CENTER)

# =============================================
# SLIDE 11 - VANTAGGI ECONOMICI
# =============================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BIANCO)
add_shape_bg(slide, 0, 0, 13.333, 1.2, BLU)
add_shape_bg(slide, 0, 1.15, 13.333, 0.05, ORO)
add_textbox(slide, 0.8, 0.25, 11.733, 0.7, 'Vantaggi economici', 32, BIANCO, True)

# Colonna sinistra - costi vecchio sito
add_shape_bg(slide, 0.6, 1.5, 5.9, 4.5, RGBColor(0xFD, 0xE8, 0xE8))
add_textbox(slide, 0.8, 1.6, 5.5, 0.5, 'Sito attuale (Joomla)', 22, ROSSO, True, PP_ALIGN.CENTER)
add_shape_bg(slide, 0.6, 2.1, 5.9, 0.04, ROSSO)

old_costs = [
    ('Hosting PHP + MySQL', '50 - 150 euro/anno'),
    ('Manutenzione CMS', 'Tempo-uomo significativo'),
    ('Licenze template/plugin', 'Variabile'),
    ('Interventi di emergenza', 'Frequenti e costosi'),
    ('', ''),
    ('TOTALE STIMATO', '100 - 300+ euro/anno'),
]

for i, (label, value) in enumerate(old_costs):
    y = 2.3 + i * 0.55
    if i == len(old_costs) - 1:
        add_shape_bg(slide, 0.8, y - 0.05, 5.5, 0.5, ROSSO)
        add_textbox(slide, 0.8, y, 2.8, 0.4, label, 15, BIANCO, True)
        add_textbox(slide, 3.6, y, 2.7, 0.4, value, 15, BIANCO, True, PP_ALIGN.RIGHT)
    elif label:
        add_textbox(slide, 0.8, y, 2.8, 0.4, label, 14, TESTO)
        add_textbox(slide, 3.6, y, 2.7, 0.4, value, 14, ROSSO, True, PP_ALIGN.RIGHT)

# Colonna destra - costi nuovo sito
add_shape_bg(slide, 6.8, 1.5, 5.9, 4.5, RGBColor(0xD4, 0xF5, 0xD4))
add_textbox(slide, 7.0, 1.6, 5.5, 0.5, 'Nuovo sito (Statico)', 22, VERDE, True, PP_ALIGN.CENTER)
add_shape_bg(slide, 6.8, 2.1, 5.9, 0.04, VERDE)

new_costs = [
    ('Hosting', 'GRATUITO'),
    ('Manutenzione', 'ZERO'),
    ('Licenze software', 'Nessuna (open source)'),
    ('Interventi emergenza', 'Rari e semplici'),
    ('', ''),
    ('TOTALE', '0 euro/anno'),
]

for i, (label, value) in enumerate(new_costs):
    y = 2.3 + i * 0.55
    if i == len(new_costs) - 1:
        add_shape_bg(slide, 7.0, y - 0.05, 5.5, 0.5, VERDE)
        add_textbox(slide, 7.0, y, 2.8, 0.4, label, 15, BIANCO, True)
        add_textbox(slide, 9.8, y, 2.7, 0.4, value, 15, BIANCO, True, PP_ALIGN.RIGHT)
    elif label:
        add_textbox(slide, 7.0, y, 2.8, 0.4, label, 14, TESTO)
        add_textbox(slide, 9.8, y, 2.7, 0.4, value, 14, VERDE, True, PP_ALIGN.RIGHT)

# Risparmio
add_shape_bg(slide, 2.5, 6.3, 8.333, 0.9, BLU)
add_textbox(slide, 2.7, 6.4, 8, 0.7,
    'Risparmio annuo: 100 - 300+ euro + ore di manutenzione risparmiate',
    22, ORO, True, PP_ALIGN.CENTER)

# =============================================
# SLIDE 10 - CONCLUSIONI
# =============================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_shape_bg(slide, 0, 0, 13.333, 7.5, BLU_SCURO)
add_shape_bg(slide, 0, 3.0, 13.333, 0.05, ORO)

add_textbox(slide, 1, 0.8, 11.333, 0.8, 'Conclusioni', 40, BIANCO, True, PP_ALIGN.CENTER)

conclusions = [
    '\u2713  Esperienza utente superiore per famiglie, studenti e personale',
    '\u2713  Conformità alle Linee Guida AgID e WCAG 2.1 AA',
    '\u2713  Azzeramento dei costi di hosting e manutenzione',
    '\u2713  Sicurezza intrinseca senza aggiornamenti continui',
    '\u2713  Autonomia gestionale per il team digitale',
    '\u2713  Accessibilità certificabile (Dichiarazione su form.agid.gov.it)',
]

for i, item in enumerate(conclusions):
    add_textbox(slide, 1.5, 1.8 + i * 0.55, 10.333, 0.5, item, 20, BIANCO)

add_textbox(slide, 1, 3.5, 11.333, 1.2,
    'Si sottopone la presente proposta all\'approvazione del Team Digitale\nper procedere alla pubblicazione del nuovo sito istituzionale.',
    22, ORO, True, PP_ALIGN.CENTER)

add_textbox(slide, 1, 5.5, 11.333, 0.8,
    'IISS "Giudici Saetta e Livatino"\nRavanusa e Campobello di Licata (AG)',
    18, RGBColor(0xAA, 0xAA, 0xCC), False, PP_ALIGN.CENTER)

add_textbox(slide, 1, 6.5, 11.333, 0.5,
    'Febbraio 2026', 16, RGBColor(0x88, 0x88, 0xAA), False, PP_ALIGN.CENTER)

# =============================================
# SLIDE FINALE - FRASE AD EFFETTO
# =============================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_shape_bg(slide, 0, 0, 13.333, 7.5, BLU_SCURO)

# Linee decorative oro
add_shape_bg(slide, 2, 2.0, 9.333, 0.03, ORO)
add_shape_bg(slide, 2, 5.3, 9.333, 0.03, ORO)

add_textbox(slide, 1, 2.3, 11.333, 2.5,
    '"Non si tratta di cambiare un sito web.\nSi tratta di dare alla nostra comunità scolastica\nlo strumento digitale che merita."',
    30, BIANCO, True, PP_ALIGN.CENTER)

add_textbox(slide, 1, 5.6, 11.333, 0.6,
    'Prof. Enrico Maria Caruso',
    20, ORO, True, PP_ALIGN.CENTER)
add_textbox(slide, 1, 6.1, 11.333, 0.5,
    'Docente di Informatica - IISS "Giudici Saetta e Livatino"',
    16, RGBColor(0xBB, 0xBB, 0xDD), False, PP_ALIGN.CENTER)

# =============================================
# SALVA
# =============================================
output_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'Presentazione_Nuovo_Sito_v3.pptx')
prs.save(output_path)
print(f'Presentazione salvata in: {output_path}')
